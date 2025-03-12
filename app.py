from flask import Flask, request, jsonify, send_file, redirect
from werkzeug.utils import secure_filename
from pptx import Presentation
from PIL import Image
import os
import io
import uuid

app = Flask(__name__, static_folder='static', static_url_path='')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['PROCESSED_FOLDER'] = 'static/processed'

# Ensure upload and processed directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pptx', 'ppt'}

def process_presentation(file_path):
    prs = Presentation(file_path)
    image_paths = []
    
    for i, slide in enumerate(prs.slides):
        # Create a unique filename for this slide
        image_filename = f"{uuid.uuid4()}_slide_{i+1}.png"
        image_path = os.path.join(app.config['PROCESSED_FOLDER'], image_filename)
        
        # Save slide as image
        with open(image_path, 'wb') as f:
            image = slide.shapes.title
            # Add actual slide to PNG conversion here
            # This is a placeholder - you'll need to implement actual conversion
            img = Image.new('RGB', (800, 600), 'white')
            img.save(image_path)
        
        image_paths.append(f"/processed/{image_filename}")
    
    return image_paths

@app.route('/')
def index():
    return app.send_static_file('index.html')

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        try:
            image_paths = process_presentation(file_path)
            # Clean up the uploaded file
            os.remove(file_path)
            return jsonify({'slides': image_paths})
        except Exception as e:
            return jsonify({'error': str(e)}), 500
    
    return jsonify({'error': 'Invalid file type'}), 400

# Catch-all route to handle all other routes and redirect to index
@app.route('/<path:path>')
def catch_all(path):
    return redirect('/')

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000) 